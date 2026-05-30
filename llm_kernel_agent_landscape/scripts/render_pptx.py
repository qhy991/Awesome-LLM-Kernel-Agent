#!/usr/bin/env python3
"""Render an editable PowerPoint version of the report landscape.

Usage:
    python scripts/render_pptx.py
"""
from __future__ import annotations

from collections import defaultdict
from pathlib import Path
from typing import Any

import yaml

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover
    raise SystemExit("python-pptx is required. Install with: pip install python-pptx") from exc

ROOT = Path(__file__).resolve().parents[1]
DATA = ROOT / "data"
OUT = ROOT / "figures"

SLIDE_W = 16
SLIDE_H = 9

# Coordinate conversion from SVG canvas to ppt inches.
SVG_W = 1800
SVG_H = 1050


def sx(x: float) -> float:
    return x / SVG_W * SLIDE_W


def sy(y: float) -> float:
    return y / SVG_H * SLIDE_H


def load_yaml(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_text_box(slide, text, left, top, width, height, font_size=14, color="#0B1F3A", bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = hex_to_rgb(color)
    return shape


def add_round_rect(slide, left, top, width, height, fill, line, text="", font_size=13, bold=True, text_color=None, line_width=1.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill)
    shape.line.color.rgb = hex_to_rgb(line)
    shape.line.width = Pt(line_width)
    if text:
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = hex_to_rgb(text_color or line)
    return shape


def compute_positions(entries: list[dict[str, Any]], cfg: dict[str, Any]) -> dict[str, tuple[float, float]]:
    periods = cfg["periods"]
    categories = cfg["categories"]
    visible = [e for e in entries if e.get("show_in_report", True)]
    groups: dict[tuple[str, str], list[dict[str, Any]]] = defaultdict(list)
    for e in visible:
        groups[(e["category"], e["period"])].append(e)
    positions: dict[str, tuple[float, float]] = {}
    for (cat, period), group in groups.items():
        x = categories[cat]["x"]
        y0 = periods[period]["y"]
        group = sorted(group, key=lambda e: (e.get("display_priority", 99), e["name"]))
        n = len(group)
        spacing = 54
        for i, e in enumerate(group):
            y = y0 + (i - (n - 1) / 2.0) * spacing + e.get("y_offset", 0)
            positions[e["id"]] = (x + 22 + e.get("x_offset", 0), y)
    return positions


def main() -> None:
    cfg = load_yaml(DATA / "categories.yaml")
    entries = load_yaml(DATA / "entries.yaml")
    categories = cfg["categories"]
    periods = cfg["periods"]

    OUT.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # title
    add_text_box(slide, "LLM-driven Kernel Engineering Landscape", 0, 0.18, 16, 0.5, 26, bold=True)
    add_text_box(
        slide,
        "Model capability · Agent workflow · Data · Evaluation · Systems",
        0,
        0.68,
        16,
        0.3,
        11,
        color="#334155",
        bold=True,
    )
    add_text_box(slide, "Time / maturity", 0.12, 1.38, 1.1, 0.3, 10, color="#334155", bold=True)

    # year bands
    for period, meta in periods.items():
        y = sy(meta["y"])
        # dashed line approximate as a thin gray line
        line = slide.shapes.add_connector(1, Inches(sx(130)), Inches(y), Inches(sx(1660)), Inches(y))
        line.line.color.rgb = RGBColor(170, 177, 189)
        line.line.width = Pt(0.7)
        add_round_rect(slide, sx(24), y - 0.18, sx(88), 0.36, "#F8FAFC", "#CBD5E1", period, 11, True, "#334155", 0.9)

    # headers and vertical guides
    for cat, meta in categories.items():
        x = sx(meta["x"])
        color = meta["color"]
        add_round_rect(slide, x - 1.05, sy(135), 2.1, 0.50, "#FFFFFF", color, cat, 16, True, color, 1.2)
        add_text_box(slide, meta.get("title_en", meta.get("title_cn", "")), x - 1.05, sy(182), 2.1, 0.16, 5.5, color=color, bold=True)
        guide = slide.shapes.add_connector(1, Inches(sx(meta["x"] - 90)), Inches(sy(208)), Inches(sx(meta["x"] - 90)), Inches(sy(805)))
        guide.line.color.rgb = hex_to_rgb(color)
        guide.line.width = Pt(1.2)
        for pmeta in periods.values():
            add_round_rect(slide, sx(meta["x"] - 98), sy(pmeta["y"] - 8), sx(16), sy(16), "#FFFFFF", color, "", 1, False, color, 1)

    positions = compute_positions(entries, cfg)
    for e in entries:
        if not e.get("show_in_report", True):
            continue
        meta = categories[e["category"]]
        color = meta["color"]
        fill = meta.get("soft_fill", "#FFFFFF")
        x, y = positions[e["id"]]
        line_width = 2.6 if e.get("highlight", False) else 1.1
        line_color = "#F6A800" if e.get("highlight", False) else color
        label = e["name"]
        aff = e.get("affiliation_short", "")
        if aff:
            label = f"{e['name']}\n{aff}"
        add_round_rect(slide, sx(x - 90), sy(y - 32), sx(180), sy(64), fill, line_color, label, 9.5, True, color, line_width)

    # legend
    add_round_rect(slide, sx(1460), sy(850), sx(250), sy(110), "#FFFFFF", "#CBD5E1", "", 1, False, "#334155", 1)
    add_text_box(
        slide,
        "Column = research area\nLine = time band\nGold border = highlight",
        sx(1490),
        sy(870),
        sx(190),
        sy(80),
        8.5,
        color="#334155",
        align=PP_ALIGN.LEFT,
    )

    prs.save(OUT / "landscape_report_editable.pptx")
    print(f"Wrote {OUT / 'landscape_report_editable.pptx'}")


if __name__ == "__main__":
    main()
